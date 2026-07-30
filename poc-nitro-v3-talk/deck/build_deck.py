#!/usr/bin/env python3
"""
Generates 'Nitro-v3-5min.pptx' — a 9-slide, 5-minute talk with speaker notes.
The demo is editor-driven; the file-by-file script lives in TOUR.md.

All numbers in this deck were measured on the POC in this repo, not copied
from marketing material. See RUNBOOK.md for how they were produced.

Usage: python3 deck/build_deck.py
"""

import re
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

# ---------------------------------------------------------------- palette
BG        = RGBColor(0x0B, 0x0E, 0x14)
PANEL     = RGBColor(0x14, 0x19, 0x22)
LINE      = RGBColor(0x2A, 0x33, 0x43)
FG        = RGBColor(0xE6, 0xED, 0xF3)
DIM       = RGBColor(0x8B, 0x98, 0xA9)
ACCENT    = RGBColor(0xF2, 0xA0, 0x3D)   # nitro orange
GREEN     = RGBColor(0x4A, 0xDE, 0x80)
RED       = RGBColor(0xF8, 0x71, 0x71)
PURPLE    = RGBColor(0xC7, 0x92, 0xEA)
YELLOW    = RGBColor(0xFF, 0xCB, 0x6B)
ORANGE2   = RGBColor(0xF7, 0x8C, 0x6C)

BODY_FONT = "Calibri"
MONO_FONT = "Menlo"

W, H = Inches(13.333), Inches(7.5)

prs = Presentation()
prs.slide_width, prs.slide_height = W, H
BLANK = prs.slide_layouts[6]


# ---------------------------------------------------------------- helpers
def slide(notes=""):
    s = prs.slides.add_slide(BLANK)
    bg = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, W, H)
    bg.fill.solid()
    bg.fill.fore_color.rgb = BG
    bg.line.fill.background()
    bg.shadow.inherit = False
    if notes:
        s.notes_slide.notes_text_frame.text = notes.strip()
    return s


def textbox(s, x, y, w, h, anchor=MSO_ANCHOR.TOP):
    tb = s.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
    return tf


def para(tf, first=False):
    return tf.paragraphs[0] if first else tf.add_paragraph()


def run(p, text, size=18, color=FG, bold=False, font=BODY_FONT, italic=False):
    r = p.add_run()
    r.text = text
    r.font.size = Pt(size)
    r.font.color.rgb = color
    r.font.bold = bold
    r.font.italic = italic
    r.font.name = font
    return r


def eyebrow(s, text):
    tf = textbox(s, Inches(0.85), Inches(0.55), Inches(11.6), Inches(0.4))
    p = para(tf, first=True)
    r = run(p, text.upper(), size=12, color=ACCENT, bold=True)
    r.font.name = BODY_FONT
    p.space_after = Pt(0)


def title(s, text, y=Inches(0.95), size=40):
    tf = textbox(s, Inches(0.85), y, Inches(11.6), Inches(1.0))
    p = para(tf, first=True)
    run(p, text, size=size, color=FG, bold=True)
    p.space_after = Pt(0)


def subtitle(s, text, y=Inches(1.85), size=17, color=DIM):
    tf = textbox(s, Inches(0.85), y, Inches(11.6), Inches(0.6))
    p = para(tf, first=True)
    run(p, text, size=size, color=color)
    p.space_after = Pt(0)


def panel(s, x, y, w, h, fill=PANEL, edge=LINE, accent=None):
    """Rounded panel. `accent` draws a 4px left rule."""
    sh = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, w, h)
    sh.fill.solid()
    sh.fill.fore_color.rgb = fill
    sh.line.color.rgb = edge
    sh.line.width = Pt(1)
    sh.shadow.inherit = False
    sh.adjustments[0] = 0.04
    if accent:
        bar = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, x, y + Inches(0.1),
                                 Pt(3.5), h - Inches(0.2))
        bar.fill.solid()
        bar.fill.fore_color.rgb = accent
        bar.line.fill.background()
        bar.shadow.inherit = False
    return sh


def hang(p, marl=26, ind=-26):
    """Hanging indent so wrapped lines align past the bullet glyph (points)."""
    pPr = p._p.get_or_add_pPr()
    pPr.set("marL", str(int(Pt(marl))))
    pPr.set("indent", str(int(Pt(ind))))


def bullets(s, items, x=Inches(0.85), y=Inches(2.5), w=Inches(11.6),
            size=19, gap=14):
    """items: list of (bold_lead, rest) or plain str."""
    tf = textbox(s, x, y, w, Inches(4.2))
    for i, it in enumerate(items):
        p = para(tf, first=(i == 0))
        p.space_after = Pt(gap)
        hang(p, 22, -22)
        run(p, "▸  ", size=size, color=ACCENT, bold=True)
        if isinstance(it, tuple):
            run(p, it[0], size=size, color=FG, bold=True)
            run(p, it[1], size=size, color=DIM)
        else:
            run(p, it, size=size, color=FG)


# ------------------------------------------------------- code highlighting
KEYWORDS = {
    "import", "from", "export", "default", "const", "let", "await", "async",
    "return", "new", "throw", "function", "if", "as", "type", "interface",
    "true", "false", "null", "catch",
}
NITRO_API = {
    "defineHandler", "defineCachedFunction", "defineCachedHandler",
    "defineConfig", "useStorage", "useDatabase", "HTTPError", "Response",
    "Request", "fetch", "performance",
}
TOKEN_RE = re.compile(
    r'(//[^\n]*)'                      # 1 comment
    r'|("[^"]*"|\'[^\']*\'|`[^`]*`)'   # 2 string
    r'|(\b\d+(?:\.\d+)?\b)'            # 3 number
    r'|(\b[A-Za-z_$][\w$]*\b)'         # 4 word
)


def code_panel(s, snippet, x, y, w, h, size=13.5, caption=None):
    panel(s, x, y, w, h, fill=RGBColor(0x0E, 0x13, 0x1B))
    if caption:
        ctf = textbox(s, x + Inches(0.28), y + Inches(0.18),
                      w - Inches(0.5), Inches(0.3))
        cp = para(ctf, first=True)
        run(cp, caption, size=11.5, color=ACCENT, bold=True, font=MONO_FONT)
        cp.space_after = Pt(0)
        top = y + Inches(0.62)
    else:
        top = y + Inches(0.26)

    tf = textbox(s, x + Inches(0.28), top, w - Inches(0.5), h - Inches(0.5))
    for i, line in enumerate(snippet.strip("\n").split("\n")):
        p = para(tf, first=(i == 0))
        p.space_after = Pt(1.5)
        p.line_spacing = 1.06
        if not line.strip():
            run(p, " ", size=size, font=MONO_FONT)
            continue
        indent = len(line) - len(line.lstrip())
        if indent:
            run(p, " " * indent, size=size, font=MONO_FONT)
        pos, body = 0, line.lstrip()
        for m in TOKEN_RE.finditer(body):
            if m.start() > pos:
                run(p, body[pos:m.start()], size=size, color=FG, font=MONO_FONT)
            comment, string, num, word = m.groups()
            if comment:
                run(p, comment, size=size, color=DIM, font=MONO_FONT, italic=True)
            elif string:
                run(p, string, size=size, color=GREEN, font=MONO_FONT)
            elif num:
                run(p, num, size=size, color=ORANGE2, font=MONO_FONT)
            else:
                if word in KEYWORDS:
                    run(p, word, size=size, color=PURPLE, font=MONO_FONT)
                elif word in NITRO_API:
                    run(p, word, size=size, color=YELLOW, font=MONO_FONT)
                else:
                    run(p, word, size=size, color=FG, font=MONO_FONT)
            pos = m.end()
        if pos < len(body):
            run(p, body[pos:], size=size, color=FG, font=MONO_FONT)


def table(s, headers, rows, x, y, w, col_w=None, row_h=Inches(0.42),
          size=15, highlight_col=None):
    n_cols = len(headers)
    gfx = s.shapes.add_table(len(rows) + 1, n_cols, x, y, w,
                             row_h * (len(rows) + 1))
    tbl = gfx.table
    tbl.first_row = True
    if col_w:
        for i, cw in enumerate(col_w):
            tbl.columns[i].width = cw
    for i in range(len(rows) + 1):
        tbl.rows[i].height = row_h

    for c, htext in enumerate(headers):
        cell = tbl.cell(0, c)
        cell.fill.solid()
        cell.fill.fore_color.rgb = RGBColor(0x1C, 0x23, 0x2F)
        cell.vertical_anchor = MSO_ANCHOR.MIDDLE
        cell.margin_left = cell.margin_right = Inches(0.13)
        p = cell.text_frame.paragraphs[0]
        p.alignment = PP_ALIGN.LEFT if c == 0 else PP_ALIGN.RIGHT
        run(p, htext, size=size - 1.5, color=ACCENT, bold=True)

    for r, rowvals in enumerate(rows, start=1):
        for c, val in enumerate(rowvals):
            cell = tbl.cell(r, c)
            cell.fill.solid()
            cell.fill.fore_color.rgb = (RGBColor(0x11, 0x16, 0x1E) if r % 2
                                        else RGBColor(0x15, 0x1B, 0x25))
            cell.vertical_anchor = MSO_ANCHOR.MIDDLE
            cell.margin_left = cell.margin_right = Inches(0.13)
            p = cell.text_frame.paragraphs[0]
            p.alignment = PP_ALIGN.LEFT if c == 0 else PP_ALIGN.RIGHT
            color = FG
            bold = False
            font = BODY_FONT
            if c == 0:
                font = MONO_FONT
                color = FG
            if highlight_col is not None and c == highlight_col:
                color, bold = GREEN, True
            if val.startswith("✓"):
                color, bold = GREEN, True
            elif val.startswith("✗"):
                color, bold = RED, True
            run(p, val, size=size - 1.5, color=color, bold=bold, font=font)
    return tbl


def stat_row(s, stats, y=Inches(5.55), x=Inches(0.85), total_w=Inches(11.6)):
    """stats: list of (big, label, color)"""
    n = len(stats)
    gap = Inches(0.22)
    cw = int((total_w - gap * (n - 1)) / n)
    for i, (big, label, color) in enumerate(stats):
        cx = x + i * (cw + gap)
        panel(s, cx, y, Emu(cw), Inches(1.25), accent=color)
        tf = textbox(s, cx + Inches(0.3), y + Inches(0.2),
                     Emu(cw) - Inches(0.45), Inches(0.9))
        p = para(tf, first=True)
        run(p, big, size=25, color=color, bold=True)
        p.space_after = Pt(2)
        p2 = para(tf)
        run(p2, label, size=12, color=DIM)


def footer(s, text):
    tf = textbox(s, Inches(0.85), Inches(6.92), Inches(11.6), Inches(0.35))
    p = para(tf, first=True)
    run(p, text, size=11.5, color=RGBColor(0x5A, 0x67, 0x78), font=MONO_FONT)


# ================================================================ SLIDE 1
s = slide("""
[0:00 — 15s]

Nitro. Most people meet it as "the thing under Nuxt". The interesting part is
that it stopped being a Nuxt thing.

Five minutes: what it is, how it compiles, what it costs, and a live demo.

— — — if asked: v3 went beta this year; that's the version I built the POC on.
""")
bar = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, Pt(6), H)
bar.fill.solid(); bar.fill.fore_color.rgb = ACCENT
bar.line.fill.background(); bar.shadow.inherit = False

tf = textbox(s, Inches(1.1), Inches(2.15), Inches(11.0), Inches(0.5))
p = para(tf, first=True)
run(p, "SERVER TOOLKIT  ·  UNJS  ·  v3 BETA", size=13, color=ACCENT, bold=True)

tf = textbox(s, Inches(1.1), Inches(2.75), Inches(11.0), Inches(1.4))
p = para(tf, first=True)
run(p, "Nitro", size=76, color=FG, bold=True)
run(p, "  v3", size=76, color=ACCENT, bold=True)

tf = textbox(s, Inches(1.1), Inches(4.2), Inches(10.5), Inches(1.0))
p = para(tf, first=True)
run(p, "Write your server once.\nLet the build target the platform.",
    size=25, color=DIM)

tf = textbox(s, Inches(1.1), Inches(6.35), Inches(11.0), Inches(0.5))
p = para(tf, first=True)
run(p, "Mateus  ·  5-minute talk  ·  live demo included",
    size=14, color=RGBColor(0x5A, 0x67, 0x78), font=MONO_FONT)

# ================================================================ SLIDE 2
s = slide("""
[0:15 — 25s]

Nitro is a server toolkit, not a framework. A framework owns your request
handler. Nitro owns your build and your deploy, and leaves the handler to you.

In v3 a handler is plain web standard — Request in, Response out. That's why
it isn't just Nuxt anymore: TanStack Start, Analog and SolidStart all sit on it.

— — — if asked: you can bring Hono or Elysia as the entry and still keep all
the deployment machinery. Nothing proprietary to learn or unlearn.
""")
eyebrow(s, "What it is")
title(s, "A build tool for servers")
subtitle(s, "It does not want to own your request handler. It wants to own your deploy.")
bullets(s, [
    ("Web standards, not abstractions  ", "— Request, Response, Headers, URL. h3 v2 was rewritten around them."),
    ("Filesystem routing  ", "— compiled at build time, so there is no runtime router to pay for."),
    ("Runtime-agnostic  ", "— one codebase for Node, Bun, Deno, Workers, Lambda."),
    ("Bring your own framework  ", "— Hono or Elysia as the entry, Nitro still does the build."),
], y=Inches(2.65))
stat_row(s, [
    ("Nuxt 5", "ships Nitro v3 at its core", ACCENT),
    ("TanStack Start", "built on Nitro", ACCENT),
    ("Analog · SolidStart", "built on Nitro", ACCENT),
])
footer(s, "unjs ecosystem: h3 · srvx · unstorage · ocache · db0 · unenv · rolldown")

# ================================================================ SLIDE 3
s = slide("""
[0:40 — 30s]

The mental model is a compiler.

Your routes and config go in. h3 handles routing, srvx normalises the HTTP
server across runtimes, unenv polyfills the Node APIs the edge doesn't have,
and Rolldown — a Rust bundler — builds it. Out comes one self-contained folder
shaped for one platform.

Key word: compiled. Routes are resolved at build time and code-split, so
hitting slash api slash users never loads the code for posts.

— — — if asked: batteries are opt-in and tree-shaken. Don't import the SQL
layer and it isn't in your bundle. That's how 25 deps stays 25 deps.
""")
eyebrow(s, "How it works")
title(s, "It compiles your server")

boxes_y = Inches(2.6)
bh = Inches(1.5)
cols = [
    ("YOUR CODE", "server/api/*\nserver/routes/*\nnitro.config.ts", FG),
    ("NITRO CORE", "h3 v2 — routing\nsrvx — HTTP\nunenv — polyfills", ACCENT),
    ("ROLLDOWN", "Rust bundler\ntree-shake\ncode-split", PURPLE),
    (".OUTPUT/", "one folder,\nself-contained,\nplatform-shaped", GREEN),
]
bw = Inches(2.55)
gap_x = Inches(0.42)
for i, (head, body, col) in enumerate(cols):
    x = Inches(0.85) + i * (bw + gap_x)
    panel(s, x, boxes_y, bw, bh, accent=col)
    tf = textbox(s, x + Inches(0.28), boxes_y + Inches(0.2),
                 bw - Inches(0.45), Inches(1.2))
    p = para(tf, first=True)
    run(p, head, size=13, color=col, bold=True, font=MONO_FONT)
    p.space_after = Pt(7)
    for ln in body.split("\n"):
        pp = para(tf)
        pp.space_after = Pt(1)
        run(pp, ln, size=12.5, color=DIM, font=MONO_FONT)
    if i < len(cols) - 1:
        atf = textbox(s, x + bw + Inches(0.04), boxes_y + Inches(0.55),
                      gap_x, Inches(0.4))
        ap = para(atf, first=True)
        ap.alignment = PP_ALIGN.CENTER
        run(ap, "→", size=22, color=LINE, bold=True)

panel(s, Inches(0.85), Inches(4.55), Inches(11.6), Inches(1.05), accent=YELLOW)
tf = textbox(s, Inches(1.2), Inches(4.75), Inches(11.0), Inches(0.7))
p = para(tf, first=True)
run(p, "Batteries are opt-in and tree-shaken: ", size=17, color=FG, bold=True)
run(p, "KV storage, SWR caching, SQL, WebSockets, scheduled tasks.",
    size=17, color=DIM)
p.space_after = Pt(4)
p2 = para(tf)
run(p2, "Don't import it and it is not in your bundle.", size=15, color=DIM, italic=True)

stat_row(s, [
    ("25", "transitive deps for nitro v3 (measured)", GREEN),
    ("273", "transitive deps for nitropack v2 (measured)", RED),
    ("22 ms", "dev rebuild in my POC", ACCENT),
], y=Inches(5.85))

# ================================================================ SLIDE 4
s = slide("""
[1:10 — 25s]

This is the payoff, and these are numbers from my POC last night — not their
marketing page.

One command. NITRO_PRESET picks the target. Same source, five outputs.

Notice Cloudflare is smaller than Node — unenv stripped the Node compatibility
layer it doesn't need there. Each target gets its own build.

— — — if asked: ~20 presets, and it auto-detects the platform in CI, so in
practice you often set nothing at all.
""")
eyebrow(s, "The core idea — presets")
title(s, "One codebase, five builds")
subtitle(s, "Measured on the POC in this repo. Same source files, nothing changed between runs.")
code_panel(s, "NITRO_PRESET=cloudflare_module npx nitro build",
           Inches(0.85), Inches(2.45), Inches(11.6), Inches(0.72), size=15)
table(s,
      ["preset", "output size", "gzipped", "runs on"],
      [
          ["node (default)", "114 kB", "33.3 kB", "✓ verified locally"],
          ["cloudflare_module", "88.2 kB", "26.6 kB", "✓ build passes"],
          ["vercel", "87.5 kB", "26.3 kB", "✓ build passes"],
          ["deno_deploy", "95.7 kB", "29.0 kB", "✓ build passes"],
          ["bun", "96.2 kB", "29.3 kB", "✓ build passes"],
      ],
      Inches(0.85), Inches(3.45), Inches(11.6),
      col_w=[Inches(4.0), Inches(2.5), Inches(2.3), Inches(2.8)],
      row_h=Inches(0.44), size=16)
footer(s, "~20 presets total  ·  auto-detected in CI  ·  each handler is its own code-split chunk (359 B for /api/hello)")

# ================================================================ SLIDE 5
s = slide("""
[1:35 — 22s]   (you'll open the real files in the demo, so keep this brisk)

Two ideas, then I'll show you the real thing.

Left: the filename is the route and the method — no registration file, so no
drift. Right: defineCachedFunction wraps any async function in a
stale-while-revalidate cache. Notice what's missing — no Redis client, no
connection string.

— — — if asked: it's backed by the storage layer, so dev writes files to disk
and production points at Redis or Workers KV by changing config, never this file.
""")
eyebrow(s, "Code")
title(s, "Routing is the filename. Caching is a wrapper.")
code_panel(s, """
// server/api/hello.get.ts  ->  GET /api/hello
import { defineHandler } from "nitro";

export default defineHandler((event) => {
  // event.req is a real web Request
  return { hello: "Nitro",
           method: event.req.method };
});

// hello.post.ts -> POST, same path.
// No router file. Nothing to keep in sync.
""", Inches(0.85), Inches(2.35), Inches(5.65), Inches(4.15),
    size=13, caption="filesystem routing")

code_panel(s, """
// server/api/stars/[...repo].ts
const cachedStars = defineCachedFunction(
  async (repo: string) => {
    const r = await fetch(gh + repo);
    return (await r.json()).stargazers_count;
  },
  { name: "ghStars",
    maxAge: 60,       // fresh for 60s
    staleMaxAge: -1,  // then stale + bg
    getKey: (repo) => repo },
);

// No redis client. No conn string.
""", Inches(6.8), Inches(2.35), Inches(5.65), Inches(4.15),
    size=13, caption="stale-while-revalidate in 10 lines")
footer(s, "[...repo] is a catch-all param  ·  /api/stars/nitrojs/nitro  ->  repo = \"nitrojs/nitro\"")

# ================================================================ SLIDE 7
s = slide("""
[1:57 — 110s]  >>> SWITCH TO EDITOR <<<   full script: TOUR.md

START ./demo.sh build-all IN A SECOND TERMINAL NOW — it takes 30-40s.

0. Sidebar. Point at server/. "Seven files. That's the entire server —
   routing, caching, auth, storage." Cmd+B to hide the sidebar.

1. OPEN server/api/stars/[...repo].ts   (35s)
   "The filename is the route. [...repo] is a catch-all, so /api/stars/
   nitrojs/nitro arrives as one string."
   Point at defineCachedFunction: "This wrapper is the whole caching story.
   Look at what is NOT here: no redis client, no connection string, no
   invalidation logic."
   -> BROWSER localhost:3100, click Fetch stars. Orange ~250-360 ms NETWORK.
      Click again: green, under 1 ms, CACHE. Then hit x10 — a wall of green.
   SAY the numbers actually on your screen: "Three hundred milliseconds to zero
   point four. Same handler. The only thing I added was that wrapper."
   (Cold latency is network-dependent — I've measured 123 to 359 ms. Read it live.)

2. OPEN nitro.config.ts   (30s)
   Point at routeRules: "This is not application code. It's config. CORS,
   caching, a proxy, and basic auth, declared against URL patterns. Nitro
   compiles these into CDN headers on Vercel, _headers on Netlify, runtime
   middleware on Node."
   -> TERMINAL: ./demo.sh auth        401, then the page.
   -> OPEN server/routes/admin/index.ts
   SAY: "Zero lines of auth code in the handler. It's four words in a config
   file."

3. THE CLOSER — the terminal you started earlier   (25s)
   Five presets, five sizes. "Same source files. One environment variable
   changed. Cloudflare is smaller than Node because unenv stripped the Node
   compatibility layer."

IF IT BREAKS: don't debug on stage. Slide 5 has the same two snippets; read the
numbers off slide 4. No network kills beat 1 only. See RUNBOOK.md.
""")
eyebrow(s, "Demo")
title(s, "Live: the whole server is seven files", size=38)
panel(s, Inches(0.85), Inches(2.3), Inches(11.6), Inches(2.95), accent=ACCENT)

tf = textbox(s, Inches(1.3), Inches(2.6), Inches(10.8), Inches(2.5))
demo_items = [
    ("1.", "stars/[...repo].ts", "no redis client  →  ~300 ms becomes 0.4 ms"),
    ("2.", "nitro.config.ts", "cors, cache, proxy, auth — as config"),
    ("3.", "routes/admin/", "password on the route, no auth code in the file"),
    ("4.", "nitro build ×5", "node · cloudflare · vercel · deno · bun"),
]
for i, (num, head, rest) in enumerate(demo_items):
    p = para(tf, first=(i == 0))
    p.space_after = Pt(14)
    hang(p, 30, -30)
    run(p, num + "  ", size=20, color=ACCENT, bold=True, font=MONO_FONT)
    run(p, head, size=20, color=FG, bold=True, font=MONO_FONT)
    run(p, "   " + rest, size=17, color=DIM)

tf = textbox(s, Inches(0.85), Inches(5.7), Inches(11.6), Inches(0.9))
p = para(tf, first=True)
run(p, "localhost:3100", size=26, color=GREEN, bold=True, font=MONO_FONT)
run(p, "     live code, live cache, live builds", size=17, color=DIM)

# ================================================================ SLIDE 8
s = slide("""
[3:47 — 30s]

Now the honest half.

v3 is beta, and two things bit me in one evening. The config wouldn't load until
I installed jiti — an optional peer dependency. And routes silently 404'd until
I set serverDir, which defaults to false.

The bigger one: deploy-anywhere covers your code, not your dependencies. My
filesystem storage driver pulled in a native binary and broke the Cloudflare and
Deno builds outright.

— — — if asked: the jiti error was just "createJiti is not a function", nothing
about jiti being missing. The storage fix was devStorage — filesystem in dev,
memory in production. Portability is a constraint you design for, not free.
""")
eyebrow(s, "Trade-offs — the honest part")
title(s, "What it costs you")

col_w2 = Inches(5.65)
panel(s, Inches(0.85), Inches(2.35), col_w2, Inches(3.95), accent=GREEN)
tf = textbox(s, Inches(1.15), Inches(2.6), col_w2 - Inches(0.6), Inches(3.5))
p = para(tf, first=True)
run(p, "WORTH IT", size=14, color=GREEN, bold=True, font=MONO_FONT)
p.space_after = Pt(13)
for t in [
    "No lock-in — platform is a build flag",
    "Cache, KV, SQL, WebSockets included",
    "Tiny dep tree: 25 vs 273 packages",
    "Sub-30 ms dev rebuilds (Rolldown)",
    "Web standards — knowledge transfers out",
]:
    pp = para(tf)
    pp.space_after = Pt(11)
    hang(pp, 20, -20)
    run(pp, "+  ", size=17, color=GREEN, bold=True, font=MONO_FONT)
    run(pp, t, size=16, color=FG)

panel(s, Inches(6.8), Inches(2.35), col_w2, Inches(3.95), accent=RED)
tf = textbox(s, Inches(7.1), Inches(2.6), col_w2 - Inches(0.6), Inches(3.5))
p = para(tf, first=True)
run(p, "WHAT BIT ME", size=14, color=RED, bold=True, font=MONO_FONT)
p.space_after = Pt(13)
for t in [
    "v3 is beta — API moves, docs lag",
    "No jiti installed → config won't load",
    "serverDir defaults to false → silent 404s",
    "Portable code ≠ portable dependencies",
    "Smaller community than Express",
]:
    pp = para(tf)
    pp.space_after = Pt(11)
    hang(pp, 20, -20)
    run(pp, "−  ", size=17, color=RED, bold=True, font=MONO_FONT)
    run(pp, t, size=16, color=FG)

footer(s, "Every item on the right is something I hit building the POC for this talk.")

# ================================================================ SLIDE 9
s = slide("""
[4:17 — 22s]

Where it sits.

Express and Fastify route requests and stop there. Hono is the fair comparison:
also web standard, also multi-runtime, genuinely leaner — but it's a router, not
a build system. Next gives you batteries, Vercel-shaped.

Nitro's niche: I want the batteries and the deploy targets, without a frontend
framework or a vendor attached.

— — — if asked: it's not either-or — Nitro will run Hono as its server entry.
Dep counts are transitive packages from a clean npm install yesterday; Hono at
1 is genuinely impressive, it just does less.
""")
eyebrow(s, "Compared to the alternatives")
title(s, "Where Nitro actually sits")
table(s,
      ["", "deps*", "multi-runtime", "deploy presets", "cache + KV built in"],
      [
          ["nitro v3", "25", "✓ Node/Bun/Deno/edge", "✓ ~20", "✓ yes"],
          ["hono", "1", "✓ Node/Bun/Deno/edge", "✗ none", "✗ bring your own"],
          ["fastify", "41", "✗ Node-centric", "✗ none", "✗ plugins"],
          ["express", "65", "✗ Node only", "✗ none", "✗ bring your own"],
          ["next.js", "~30", "~ Node + edge", "~ Vercel-shaped", "✓ yes"],
      ],
      Inches(0.85), Inches(2.5), Inches(11.6),
      col_w=[Inches(2.3), Inches(1.35), Inches(3.15), Inches(2.3), Inches(2.5)],
      row_h=Inches(0.46), size=15.5)

panel(s, Inches(0.85), Inches(5.35), Inches(11.6), Inches(1.25), accent=ACCENT)
tf = textbox(s, Inches(1.2), Inches(5.58), Inches(11.0), Inches(0.9))
p = para(tf, first=True)
run(p, "Not either-or: ", size=17, color=ACCENT, bold=True)
run(p, "Nitro will run Hono, Elysia or plain h3 as its server entry — you keep the router you like and still get the presets.",
    size=16, color=FG)
p.space_after = Pt(5)
p2 = para(tf)
run(p2, "* transitive package count, measured with a clean npm install on 2026-07-29.",
    size=12, color=DIM, italic=True)

# ================================================================ SLIDE 10
s = slide("""
[4:39 — 18s]

Would I use it?

For a new API that might move platforms, yes — for the caching and storage, not
the routing. Shipping next week? Wait for stable.

The idea worth stealing: deployment as a compile target. Which platform becomes
a build flag, not an architecture decision.

Thanks.

— — — if asked: v2 (nitropack) is production-proven today and runs Nuxt 3 and 4.
""")
eyebrow(s, "Verdict")
title(s, "Deployment as a compile target")
subtitle(s, "That is the idea worth stealing, whether or not you adopt the tool.")

panel(s, Inches(0.85), Inches(2.6), Inches(5.65), Inches(2.15), accent=GREEN)
tf = textbox(s, Inches(1.15), Inches(2.85), Inches(5.05), Inches(1.7))
p = para(tf, first=True)
run(p, "REACH FOR IT WHEN", size=13, color=GREEN, bold=True, font=MONO_FONT)
p.space_after = Pt(10)
for t in ["Edge or multi-platform is on the table",
          "You want cache + KV without provisioning",
          "You're already in the Nuxt / unjs world"]:
    pp = para(tf); pp.space_after = Pt(7)
    run(pp, "·  " + t, size=15.5, color=FG)

panel(s, Inches(6.8), Inches(2.6), Inches(5.65), Inches(2.15), accent=RED)
tf = textbox(s, Inches(7.1), Inches(2.85), Inches(5.05), Inches(1.7))
p = para(tf, first=True)
run(p, "SKIP IT WHEN", size=13, color=RED, bold=True, font=MONO_FONT)
p.space_after = Pt(10)
for t in ["Production ships next week (wait for stable)",
          "A plain Node service on one host is enough",
          "Your team needs a big hiring pool / lots of docs"]:
    pp = para(tf); pp.space_after = Pt(7)
    run(pp, "·  " + t, size=15.5, color=FG)

tf = textbox(s, Inches(0.85), Inches(5.15), Inches(11.6), Inches(1.2))
p = para(tf, first=True)
run(p, "nitro.build", size=22, color=ACCENT, bold=True, font=MONO_FONT)
run(p, "     ·     ", size=22, color=LINE, font=MONO_FONT)
run(p, "github.com/nitrojs/nitro", size=22, color=ACCENT, bold=True, font=MONO_FONT)
p.space_after = Pt(10)
p2 = para(tf)
run(p2, "Demo source, measured numbers and the runbook are in this repo — ask me for it.",
    size=16, color=DIM)

tf = textbox(s, Inches(0.85), Inches(6.5), Inches(11.6), Inches(0.6))
p = para(tf, first=True)
run(p, "Questions?", size=30, color=FG, bold=True)


# ---------------------------------------------------------------- write
import os
here = os.path.dirname(os.path.abspath(__file__))
out = os.path.join(here, "Nitro-v3-5min.pptx")
prs.save(out)
print(f"Wrote {out}  ({len(prs.slides.__iter__.__self__._sldIdLst)} slides)")

# Also emit the notes as markdown, so the script survives a machine with no
# PowerPoint. Generated from the same source — it can never drift out of sync.
titles = []
for s in prs.slides:
    biggest, size = "", 0
    for sh in s.shapes:
        if not sh.has_text_frame:
            continue
        for p in sh.text_frame.paragraphs:
            for r in p.runs:
                if r.font.size and r.font.size > size and r.text.strip():
                    biggest, size = r.text.strip(), r.font.size
    titles.append(biggest)

lines = [
    "# Speaker notes — Nitro v3, 5-minute talk",
    "",
    "Gerado automaticamente por `build_deck.py`. **Não edite à mão** — edite o",
    "script e rode `python3 deck/build_deck.py`.",
    "",
    "Fallback para quando não houver PowerPoint na máquina. O roteiro do código",
    "está em [`../TOUR.md`](../TOUR.md).",
    "",
    "> Blocos `— — — if asked:` **não são pra ler** — é munição para o Q&A.",
    "",
    "---",
    "",
]
for i, s in enumerate(prs.slides, 1):
    note = s.notes_slide.notes_text_frame.text.strip()
    timing = note.split("\n")[0] if note else ""
    body = "\n".join(note.split("\n")[1:]).strip()
    lines += [f"## Slide {i} — {titles[i - 1]}", "", f"`{timing}`", "", body, "", "---", ""]

notes_md = os.path.join(here, "SPEAKER-NOTES.md")
with open(notes_md, "w") as fh:
    fh.write("\n".join(lines))
print(f"Wrote {notes_md}")
